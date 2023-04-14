import tkinter as tk

from tkinter.scrolledtext import ScrolledText
from tkinter import filedialog
from tkinter import *
from pytube import YouTube
import os
from pptx.util import Inches
import cv2
from pptx import Presentation
import threading
import re
import numpy as np
import csv

def thread_it(func, *args):
    t = threading.Thread(target=func, args=args)
    t.setDaemon(True)
    t.start()

def onProgress(stream, chunk, remains):
    total = stream.filesize # 取得完整尺寸
    percent = (total-remains)/total*100 # 減去剩餘尺寸(剩餘尺寸會抓取存取的檔案大小)
    print(f'下载中 {percent:05.2f} %', end='\r') # 顯示進度，\r 表示不換行，在同一行更新

def MSE(img1, img2):
        squared_diff = (img1-img2)**2
        summed = np.sum(squared_diff)
        num_pix = img1.shape[0]*img1.shape[1]
        err = summed/num_pix
        # print(err)
        # with open('data.csv', 'w', newline='') as file:
        #     writer = csv.writer(file)
            
        #     writer.writerow(list(err.tolist()))
        print(err)
        return err


def select_file():

    filename = filedialog.askopenfilename(initialdir = os.getcwd(), title = "Select a File",filetypes = (("Video files", "*.mp4;*.mkv;*.avi;*.wmv;*.flv"), ("all files", "*.*"))
)
    # 在entry中显示文件路径
    file_path_entry.delete(0, END)
    file_path_entry.insert(0, filename)

def download_video():
    stext.insert(tk.END, "开始下载视频...\n")
    yt = YouTube(input_url.get())
    yt.streams.filter().get_by_resolution('720p').download(filename="input.mp4")
    stext.insert(tk.END, "下载成功\n\n")



def video_to_ppt(filePath=None):
    stext.insert(tk.END, "正在生成图片，需要等待一段时间...\n")
    if not os.path.isdir("output"):
        os.mkdir('output')
    if filePath is None:
        filePath = 'input.mp4'
    os.system(f'ffmpeg -loglevel error -i {filePath} -vf fps=1 output/out%d.png')
    stext.insert(tk.END, "图片生成完成\n\n")

    stext.insert(tk.END, "开始处理图片...\n")
    #TODO 处理图片，去掉部分相似的图片

    fileNum = len([x for x in os.listdir('output/')]) # read images
    imgPrev = cv2.imread('output/out1.png') # first image
    for i in range(fileNum-1):
        filePath = 'output/out'+ str(i+2) +'.png'
        img = cv2.imread(filePath)
        grayA = cv2.cvtColor(imgPrev, cv2.COLOR_BGR2GRAY)
        grayB = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        # check threshold and delete images
        if MSE(grayA,grayB) < 2:
            os.remove(filePath)
        #print(str(i) + ' MSE: ' + str(MSE(grayA,grayB)))
        imgPrev = img

    stext.insert(tk.END, "图片处理完成\n\n")

    stext.insert(tk.END, "正在制作PowerPoint...\n")
    prs = Presentation()
    prs.slide_height = Inches(9)
    prs.slide_width = Inches(16)
    blank_slide_layout = prs.slide_layouts[0]
    #sort the images by name
    list_dir = os.listdir('output/')
    list_dir = sorted(list_dir, key=lambda s: int(re.search(r'\d+', s).group()))
    for f in list_dir:
        slide = prs.slides.add_slide(blank_slide_layout)
        pic = slide.shapes.add_picture('output/' + f, 0, 0, Inches(16), Inches(9))
        os.remove('output/' + f)
    prs.save("result.pptx")
    os.rmdir('output')
    stext.insert(tk.END, 'PowerPoint制作完成\n\n')

def end_to_end():
    download_video()
    print("end to end ...")
    video_to_ppt()



if __name__ == "__main__":
    win = tk.Tk()
    win.title("Video To Powerpoint")

    width = 400
    height = 400
    screenwidth = win.winfo_screenwidth()  
    screenheight = win.winfo_screenheight()
    print(screenwidth, screenheight) 
    alignstr = '%dx%d+%d+%d' % (width, height, (screenwidth-width)/2, (screenheight-height)/2)   
    print(alignstr)
    win.geometry(alignstr)
    win.resizable(0,0)

    instruction = tk.Label(text='输入视频网址: ', font=('楷体', 12))
    instruction.pack(padx=20, pady=10, anchor=tk.W)

    input_url = tk.Entry()
    input_url.pack(ipadx=120, ipady=6, padx=20, anchor=tk.W)
    
    frame = tk.Frame(win)
    frame.pack(side=tk.TOP)
    button_spacing =20

    button1 = tk.Button(frame, text="一键生成", font=('楷体',12), command=lambda :thread_it(end_to_end))
    button1.pack(side=tk.LEFT, padx=button_spacing, pady=10,)


    button2 = tk.Button(frame, text="下载视频", font=('楷体',12), command=lambda :thread_it(download_video))
    button2.pack(side=tk.LEFT, padx=button_spacing, pady=10)


    frame_local_generate = tk.Frame(win)
    frame_local_generate.pack(side=tk.TOP)

    spacing = 10
    # 创建entry用于显示文件路径
    file_path_entry = Entry(frame_local_generate)
    file_path_entry.pack(side=tk.LEFT, padx=spacing, ipady=6 ,pady=10)

    select_button = Button(frame_local_generate, text="选择文件", font=('楷体', 12), command=select_file)
    select_button.pack(side=tk.LEFT, padx=spacing, pady=10)

    button3 = tk.Button(frame_local_generate, text="本地生成", font=('楷体', 12),  command=lambda :thread_it(video_to_ppt, file_path_entry.get()))
    button3.pack(side=tk.LEFT, padx=spacing, pady=10)

    stext = ScrolledText(width=60, height=18, background='#F7F3EC')
    stext.pack(padx=20)

    win.mainloop()
